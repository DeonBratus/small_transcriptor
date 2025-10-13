"""
PPTX file processor with image handling.
"""

import os
import base64
import requests
from typing import List, Dict, Any
from pptx import Presentation
from .base_processor import BaseFileProcessor


class PptxProcessor(BaseFileProcessor):
    """Processor for PPTX files with image description capabilities."""
    
    # Constants for image size filtering
    EMU_PER_INCH = 914400
    CM_PER_INCH = 2.54
    
    def __init__(self, file_path: str, vision_model: str = "llava", 
                 ollama_base_url: str = "http://localhost:11434",
                 min_width_cm: float = 20, min_height_cm: float = 9,
                 images_dir: str = "/app/data/evaluation_files/extracted_images"):
        """
        Initialize PPTX processor.
        
        Args:
            file_path: Path to PPTX file
            vision_model: Model for image description
            ollama_base_url: Base URL for Ollama API
            min_width_cm: Minimum width for images to process
            min_height_cm: Minimum height for images to process
            images_dir: Directory to save extracted images
        """
        super().__init__(file_path)
        self.vision_model = vision_model
        self.ollama_base_url = ollama_base_url.rstrip('/')
        self.min_width_cm = min_width_cm
        self.min_height_cm = min_height_cm
        self.images_dir = images_dir
        self.slides_data = []
    
    def validate_file(self) -> bool:
        """Validate PPTX file."""
        if not os.path.isfile(self.file_path):
            return False
        
        # Check file extension
        if not self.file_path.lower().endswith('.pptx'):
            return False
            
        return True
    
    def describe_image(self, image_path: str) -> str:
        """
        Describe image using Ollama vision model.
        
        Args:
            image_path: Path to image file
            
        Returns:
            Description of the image
        """
        try:
            with open(image_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode("utf-8")

            # Prepare request for Ollama
            payload = {
                "model": self.vision_model,
                "prompt": "Describe this slide visual in detail, preserving any numeric data and explaining what it shows:",
                "images": [img_b64],
                "stream": False
            }

            response = requests.post(
                f"{self.ollama_base_url}/api/generate",
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get("response", "Could not describe image").strip()
            else:
                print(f"Error describing image: {response.status_code}")
                return "Image description unavailable"
                
        except Exception as e:
            print(f"Error in describe_image: {e}")
            return "Image description unavailable"
    
    def process(self) -> str:
        """
        Extract detailed content from PPTX file including images.
        
        Returns:
            Detailed transcript of the presentation
        """
        if not self.validate_file():
            raise FileNotFoundError(f"PPTX file not found or invalid: {self.file_path}")

        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            raise ValueError(f"Error opening PPTX file '{self.file_path}': {e}")

        os.makedirs(self.images_dir, exist_ok=True)
        all_slides_text = []
        self.slides_data = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_data = {
                'slide_number': slide_idx,
                'title': '',
                'text_content': [],
                'tables': [],
                'images': []
            }
            
            slide_text = [f"--- Slide {slide_idx} ---"]

            # Extract title
            title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
            first_title = title_shapes[0].text.strip() if title_shapes and title_shapes[0].text else ""
            if first_title:
                slide_text.append(f"Title: {first_title}")
                slide_data['title'] = first_title

            # Extract body text
            body_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text.strip()
                    if txt and txt != first_title:
                        body_texts.append(txt)
                        slide_data['text_content'].append(txt)
            
            if body_texts:
                slide_text.append("Text:\n" + "\n".join(f"- {t}" for t in body_texts))

            # Extract tables
            for shape in slide.shapes:
                if shape.has_table:
                    table_content = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_content.append(row_data)
                    
                    if table_content:
                        header = " | ".join(table_content[0])
                        sep = " | ".join(["---"] * len(table_content[0]))
                        rows = "\n".join(" | ".join(r) for r in table_content[1:])
                        table_md = f"| {header} |\n| {sep} |\n" + "\n".join(f"| {r} |" for r in rows.split("\n"))
                        slide_text.append("Table Content:\n" + table_md)
                        slide_data['tables'].append(table_content)

            # Extract and describe images
            for img_idx, shape in enumerate(slide.shapes, start=1):
                if shape.shape_type == 13:  # picture
                    width_cm = shape.width / self.EMU_PER_INCH * self.CM_PER_INCH
                    height_cm = shape.height / self.EMU_PER_INCH * self.CM_PER_INCH
                    
                    if width_cm >= self.min_width_cm and height_cm >= self.min_height_cm:
                        image = shape.image
                        ext = image.ext
                        img_path = os.path.join(self.images_dir, f"slide_{slide_idx}_img_{img_idx}.{ext}")
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        img_desc = self.describe_image(img_path)
                        slide_text.append(f"Image/Graph Description:\n{img_desc}")
                        
                        slide_data['images'].append({
                            'path': img_path,
                            'description': img_desc,
                            'width_cm': width_cm,
                            'height_cm': height_cm
                        })

            all_slides_text.append("\n".join(slide_text))
            self.slides_data.append(slide_data)

        self.processed_content = "\n\n".join(all_slides_text)
        return self.processed_content
    
    def get_slides_data(self) -> List[Dict[str, Any]]:
        """
        Get structured data about slides.
        
        Returns:
            List of slide data dictionaries
        """
        if not self.slides_data:
            self.process()
        return self.slides_data
    
    def get_metadata(self) -> dict:
        """Get extended metadata for PPTX files."""
        metadata = super().get_metadata()
        
        if self.slides_data:
            total_images = sum(len(slide['images']) for slide in self.slides_data)
            total_tables = sum(len(slide['tables']) for slide in self.slides_data)
            
            metadata.update({
                'slide_count': len(self.slides_data),
                'image_count': total_images,
                'table_count': total_tables,
                'vision_model': self.vision_model
            })
        
        return metadata
