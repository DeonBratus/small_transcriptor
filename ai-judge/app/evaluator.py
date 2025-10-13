import os
import base64
import requests
import json
from pptx import Presentation
import docx
from typing import Generator

EMU_PER_INCH = 914400
CM_PER_INCH = 2.54


class PresentationEvaluator:
    def __init__(self, vision_model, eval_model, ollama_base_url="http://localhost:11434", min_width_cm=20, min_height_cm=9):
        self.vision_model = vision_model
        self.eval_model = eval_model
        self.ollama_base_url = ollama_base_url.rstrip('/')
        self.min_width_cm = min_width_cm
        self.min_height_cm = min_height_cm
        self.doc_text = ''
        self.presentation_text = ''

    # --- DOCX ---
    def set_docx_transcript(self, path):
        """Load a DOCX file and extract text."""
        if not os.path.isfile(path):
            raise FileNotFoundError(f"DOCX file not found: {path}")

        try:
            doc = docx.Document(path)
            self.doc_text = "\n".join(
                p.text for p in doc.paragraphs if p.text.strip()
            )
        except Exception as e:
            raise ValueError(f"Error reading DOCX file '{path}': {e}")

    # --- Image description ---
    def describe_image(self, image_path):
        """Describe image using Ollama vision model"""
        try:
            with open(image_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode("utf-8")

            # Prepare request for Ollama
            payload = {
                "model": self.vision_model,
                "prompt": "Describe this slide visual in detail, preserving any numeric data and explaining what it shows:",
                "images": [img_b64],
                "stream": False
            }

            response = requests.post(
                f"{self.ollama_base_url}/api/generate",
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get("response", "Could not describe image").strip()
            else:
                print(f"Error describing image: {response.status_code}")
                return "Image description unavailable"
                
        except Exception as e:
            print(f"Error in describe_image: {e}")
            return "Image description unavailable"

    def set_pptx_detailed_transcript(self, path, images_dir="/app/data/evaluation_files/extracted_images"):
        """Load a PPTX file and extract a detailed transcript."""
        if not os.path.isfile(path):
            raise FileNotFoundError(f"PPTX file not found: {path}")

        try:
            prs = Presentation(path)
        except Exception as e:
            raise ValueError(f"Error opening PPTX file '{path}': {e}")

        os.makedirs(images_dir, exist_ok=True)
        all_slides_text = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text = [f"--- Slide {slide_idx} ---"]

            # Title
            title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
            first_title = title_shapes[0].text.strip() if title_shapes and title_shapes[0].text else ""
            if first_title:
                slide_text.append(f"Title: {first_title}")

            # Body text
            body_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text.strip()
                    if txt and txt != first_title:
                        body_texts.append(txt)
            if body_texts:
                slide_text.append("Text:\n" + "\n".join(f"- {t}" for t in body_texts))

            # Tables
            for shape in slide.shapes:
                if shape.has_table:
                    table_content = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_content.append(row_data)
                    if table_content:
                        header = " | ".join(table_content[0])
                        sep = " | ".join(["---"] * len(table_content[0]))
                        rows = "\n".join(" | ".join(r) for r in table_content[1:])
                        table_md = f"| {header} |\n| {sep} |\n" + "\n".join(f"| {r} |" for r in rows.split("\n"))
                        slide_text.append("Table Content:\n" + table_md)

            # Images
            for img_idx, shape in enumerate(slide.shapes, start=1):
                if shape.shape_type == 13:  # picture
                    width_cm = shape.width / EMU_PER_INCH * CM_PER_INCH
                    height_cm = shape.height / EMU_PER_INCH * CM_PER_INCH
                    if width_cm >= self.min_width_cm and height_cm >= self.min_height_cm:
                        image = shape.image
                        ext = image.ext
                        img_path = os.path.join(images_dir, f"slide_{slide_idx}_img_{img_idx}.{ext}")
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        img_desc = self.describe_image(img_path)
                        slide_text.append(f"Image/Graph Description:\n{img_desc}")

            all_slides_text.append("\n".join(slide_text))

        self.presentation_text = "\n\n".join(all_slides_text)

    # --- Model call for evaluation ---
    def evaluate(self, doc_path, presentation_path=None):
        """Evaluate document and presentation using Ollama"""
        self.set_docx_transcript(doc_path)

        if presentation_path is not None:
            self.set_pptx_detailed_transcript(presentation_path)

        if self.presentation_text:
            task_prompt = self.build_thesis_presentation_prompt(self.doc_text, self.presentation_text)
        else:
            task_prompt = self.build_thesis_only_prompt(self.doc_text)

        full_prompt = f"{self.system_prompt()}\n\n{task_prompt}"

        try:
            payload = {
                "model": self.eval_model,
                "prompt": full_prompt,
                "stream": False,
                "options": {
                    "temperature": 0.0,
                    "num_predict": 2500
                }
            }

            response = requests.post(
                f"{self.ollama_base_url}/api/generate",
                json=payload,
                timeout=300
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get("response", "Evaluation failed")
            else:
                return f"Error: {response.status_code} - {response.text}"
                
        except Exception as e:
            return f"Error during evaluation: {str(e)}"

    def evaluate_stream(self, doc_path, presentation_path=None) -> Generator[str, None, None]:
        """Stream evaluation results using Ollama"""
        self.set_docx_transcript(doc_path)

        if presentation_path is not None:
            self.set_pptx_detailed_transcript(presentation_path)

        if self.presentation_text:
            task_prompt = self.build_thesis_presentation_prompt(self.doc_text, self.presentation_text)
        else:
            task_prompt = self.build_thesis_only_prompt(self.doc_text)

        full_prompt = f"{self.system_prompt()}\n\n{task_prompt}"

        try:
            payload = {
                "model": self.eval_model,
                "prompt": full_prompt,
                "stream": True,
                "options": {
                    "temperature": 0.0,
                    "num_predict": 2500
                }
            }

            response = requests.post(
                f"{self.ollama_base_url}/api/generate",
                json=payload,
                stream=True,
                timeout=300
            )
            
            if response.status_code == 200:
                for line in response.iter_lines():
                    if line:
                        try:
                            data = json.loads(line.decode('utf-8'))
                            if 'response' in data:
                                yield data['response']
                            if data.get('done', False):
                                break
                        except json.JSONDecodeError:
                            continue
            else:
                yield f"Error: {response.status_code} - {response.text}"
                
        except Exception as e:
            yield f"Error during evaluation: {str(e)}"

    # --- Prompts ---
    @staticmethod
    def system_prompt():
        return """
You are an experienced academic reviewer and pedagogue. For every request you must produce exactly two sections in this order:

1) A short free-text reasoning section bracketed by "<THOUGHT>" and "</THOUGHT>". In this section, give concise, specific, reviewer-style notes about the submission (intuition, main strengths/weaknesses you will weigh, and reasoning steps you used). Keep this section focused and specific to the provided document (max ~200-300 words).

2) A machine-readable evaluation bracketed by "<JSON>" and "</JSON>". Output only valid JSON in this block (no extra text). Follow the schema described in the incoming task prompt (the task prompt will request either thesis-only fields or thesis+presentation fields). For each numeric rating, use the requested integer scale. For list fields, return JSON arrays of strings. For boolean fields, use true/false.

Always obey the specific fields, labels and rating scales specified in the task prompt. If the task prompt includes both a thesis and presentation, include the presentation fields. If the task prompt contains only a thesis, omit presentation fields.

Do not output anything else outside the two sections. Be concise and deterministic: set tone professional, neutral, and specific. When you mention weaknesses, be constructive and give concrete ways to improve.
"""

    @staticmethod
    def build_thesis_presentation_prompt(doc_text, presentation_text):
        return f"""
Instructions:
Evaluate the following THESIS/DOCUMENT and the PRESENTATION according to the general evaluation criteria below.
Return two sections exactly as the system prompt requires: <THOUGHT> and <JSON>.

Evaluation criteria (main, general — labels you must use in the JSON):
1. "Relevance" / "Актуальность темы доклада"
2. "Literature_and_Originality" / "Полнота обзора и оригинальность"
3. "Practical_Significance" / "Практическая значимость"
4. "Presentation_and_QA" / "Качество презентации и ответов на вопросы"
5. "Author_Contribution" / "Личный вклад автора"

--- DOCUMENT BEGIN ---
{doc_text}
--- DOCUMENT END ---

--- PRESENTATION TRANSCRIPT BEGIN ---
{presentation_text}
--- PRESENTATION TRANSCRIPT END ---

Please produce:
- <THOUGHT>: brief, specific notes about reasoning and priorities for scoring the document and presentation (max ~200-300 words).
- <JSON>: an object with the following fields (in this exact order):

- "Summary": string. Short summary of the paper and its main contributions.
- "Strengths": [string]. List of strengths.
- "Weaknesses": [string]. List of weaknesses or missing parts.
- "Relevance": integer 1-4 (low, medium, high, very high).
- "Literature_and_Originality": integer 1-4 (low, medium, high, very high).
- "Practical_Significance": integer 1-4 (low, medium, high, very high).
- "Author_Contribution": integer 1-4 (low, medium, high, very high).
- "Presentation_and_QA": integer 1-4 (low, medium, high, very high).
- "PresentationNotes": [string]. Brief, concrete notes about delivery and answers to questions.
- "Soundness": integer 1-4 (poor, fair, good, excellent).
- "Questions": [string]. Clarifying questions for the authors.
- "Limitations": [string]. Limitations and possible negative societal impacts.
- "Ethical_Concerns": boolean.
- "Overall": integer 1-10 (1 = Very Strong Reject, 2 = Strong Reject, 3 = Reject, 4 = Borderline reject, 5 = Borderline, 6 = Weak Accept, 7 = Accept, 8 = Strong Accept, 9 = Very Strong Accept, 10 = Award quality).
- "Confidence": integer 1-5 (low, medium, high, very high, absolute).
- "Decision": string, one of ["Accept", "Reject"].

Return only the two required blocks. Do not add extra commentary.
"""

    @staticmethod
    def build_thesis_only_prompt(doc_text):
        return f"""
Instructions:
Evaluate the following THESIS/DOCUMENT according to the general evaluation criteria below.
Return two sections exactly as the system prompt requires: <THOUGHT> and <JSON>.

Evaluation criteria (main, general — translated and written as labels you must use in the JSON):
1. "Relevance" / "Актуальность темы доклада" — How relevant the topic is to current challenges in the field.
2. "Literature_and_Originality" / "Полнота обзора и оригинальность" — Completeness of literature review, analysis of prior solutions and their shortcomings, and the originality/novelty / theoretical significance of the proposed solution.
3. "Practical_Significance" / "Практическая значимость" — Practical impact and the feasibility/quality of practical implementation (if applicable).
4. "Author_Contribution" / "Личный вклад автора" — Degree of author’s personal contribution to the solution.
5. "PresentationQuality" is NOT required in this variant (omit presentation fields).

--- DOCUMENT BEGIN ---
{doc_text}
--- DOCUMENT END ---

Please produce:
- <THOUGHT>: brief, specific notes about reasoning and priorities for scoring this document.
- <JSON>: an object with the following fields (in this exact order):

- "Summary": string. Short summary of the paper and its main contributions.
- "Strengths": [string]. List of strengths.
- "Weaknesses": [string]. List of weaknesses or missing parts.
- "Relevance": integer 1-4 (low, medium, high, very high).
- "Literature_and_Originality": integer 1-4 (low, medium, high, very high).
- "Practical_Significance": integer 1-4 (low, medium, high, very high).
- "Author_Contribution": integer 1-4 (low, medium, high, very high).
- "Soundness": integer 1-4 (poor, fair, good, excellent).
- "Questions": [string]. Clarifying questions for the authors.
- "Limitations": [string]. Limitations and possible negative societal impacts.
- "Ethical_Concerns": boolean.
- "Overall": integer 1-10 (1 = Very Strong Reject, 2 = Strong Reject, 3 = Reject, 4 = Borderline reject, 5 = Borderline, 6 = Weak Accept, 7 = Accept, 8 = Strong Accept, 9 = Very Strong Accept, 10 = Award quality).
- "Confidence": integer 1-5 (low, medium, high, very high, absolute).
- "Decision": string, one of ["Accept", "Reject"].

Return only the two required blocks. Do not add extra commentary.
"""


# === Example usage ===
if __name__ == "__main__":
    evaluator = PresentationEvaluator(
        vision_model="llava",
        eval_model="llama3.2",
        ollama_base_url="http://localhost:11434"
    )

    doc_path = "example.docx"
    pres_path = "example.pptx"

    result = evaluator.evaluate(doc_path, pres_path)  # testing
    print(result)
